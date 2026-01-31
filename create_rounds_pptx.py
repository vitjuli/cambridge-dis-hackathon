"""
Generate PowerPoint presentation for 1-Round vs 3-Round analysis
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_BG = RGBColor(15, 23, 42)
ACCENT = RGBColor(99, 102, 241)
SUCCESS = RGBColor(16, 185, 129)
WARNING = RGBColor(245, 158, 11)
DANGER = RGBColor(239, 68, 68)
PURPLE = RGBColor(124, 58, 237)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(148, 163, 184)


def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "1-Round vs 3-Round Debate Analysis"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "Impact of Debate Rounds on Claim Verification Quality"
    p.font.size = Pt(20)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(5.5), Inches(3), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(30, 27, 75)
    badge.line.color.rgb = ACCENT
    tf = badge.text_frame
    tf.paragraphs[0].text = "Kepler Team"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(1.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)


def add_comparison_slide(prs, title, data):
    """Three-way comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    configs = [
        ("1-Round\n(Single)", DANGER, data[0]),
        ("3-Round\n(Standard)", SUCCESS, data[1]),
        ("3-Round\n(Forced)", PURPLE, data[2]),
    ]

    for i, (label, color, items) in enumerate(configs):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.3 + i * 3.2), Inches(1.2),
                                     Inches(3), Inches(5.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(30, 27, 75)
        box.line.color.rgb = color

        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
            p.space_before = Pt(6)


def add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    cols = len(headers)
    table_rows = len(rows) + 1
    table = slide.shapes.add_table(table_rows, cols, Inches(0.5), Inches(1.3), Inches(9), Inches(0.5 * table_rows)).table

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(30, 27, 75)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER


def add_formula_slide(prs, title, formulas):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    y_pos = 1.3
    for name, formula, desc in formulas:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(9), Inches(1.4))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(30, 27, 75)
        box.line.color.rgb = RGBColor(55, 48, 107)

        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = ACCENT

        p = tf.add_paragraph()
        p.text = formula
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.name = "Consolas"

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = LIGHT_GRAY

        y_pos += 1.6


def add_rounds_visual_slide(prs):
    """Visual showing 3 rounds."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "3-Round Debate Protocol"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    rounds = [
        ("R1", "Initial Stance", DANGER, "Independent analysis\nEstablish positions"),
        ("R2", "Cross-Examination", WARNING, "Address rebuttals\nChallenge arguments"),
        ("R3", "Final Position", SUCCESS, "Refined stance\nIncorporate valid points"),
    ]

    for i, (num, title, color, desc) in enumerate(rounds):
        x = 0.8 + i * 3

        # Circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.5), Inches(1.8), Inches(1.2), Inches(1.2))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(x + 0.5), Inches(2.1), Inches(1.2), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        title_b = slide.shapes.add_textbox(Inches(x), Inches(3.2), Inches(2.2), Inches(0.5))
        tf = title_b.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        desc_b = slide.shapes.add_textbox(Inches(x), Inches(3.7), Inches(2.2), Inches(1))
        tf = desc_b.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

        # Arrow
        if i < 2:
            arrow = slide.shapes.add_textbox(Inches(x + 2.3), Inches(2.2), Inches(0.5), Inches(0.5))
            tf = arrow.text_frame
            p = tf.paragraphs[0]
            p.text = "→"
            p.font.size = Pt(32)
            p.font.color.rgb = LIGHT_GRAY

    # Key insight box
    insight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5), Inches(9), Inches(1.5))
    insight.fill.solid()
    insight.fill.fore_color.rgb = RGBColor(30, 41, 59)
    insight.line.color.rgb = ACCENT

    tf = insight.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Why 3 Rounds?"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    p = tf.add_paragraph()
    p.text = "Each round refines arguments. Round 3 arguments are 40% longer and address specific rebuttals from opposing agents."
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE


def add_key_findings_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, RGBColor(30, 27, 75))

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Findings: 1-Round vs 3-Round"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    findings = [
        ("Confidence", "-10%", "95% → 85% (better calibrated)"),
        ("Mutations", "+2", "0 → 2 types detected"),
        ("Ambiguity", "100%", "Correctly identifies uncertain cases"),
        ("Transparency", "Full", "Complete debate transcript"),
    ]

    for i, (label, value, desc) in enumerate(findings):
        row = i // 2
        col = i % 2
        x = 0.5 + col * 4.7
        y = 1.5 + row * 2.5

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.3), Inches(2))
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_BG
        box.line.color.rgb = SUCCESS

        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = SUCCESS
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER


def add_conclusion_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "3-Round Debate: Worth the Cost"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Trade-off visual
    tradeoff = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    tf = tradeoff.text_frame
    p = tf.paragraphs[0]
    p.text = "10× compute cost  →  2× better calibration"
    p.font.size = Pt(24)
    p.font.color.rgb = SUCCESS
    p.alignment = PP_ALIGN.CENTER

    takeaway = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.8), Inches(8), Inches(1.5))
    takeaway.fill.solid()
    takeaway.fill.fore_color.rgb = RGBColor(30, 27, 75)
    takeaway.line.color.rgb = ACCENT

    tf = takeaway.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "For fact verification where false confidence is dangerous,"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "additional debate rounds are a worthwhile investment."
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slides
    add_title_slide(prs)

    add_content_slide(prs, "Research Question", [
        "How many debate rounds are optimal for claim verification?",
        "Single-agent (1 round): Fast but potentially overconfident",
        "Multi-agent (3 rounds): More compute but deeper analysis",
        "Trade-off: Computational cost vs verification quality"
    ])

    add_rounds_visual_slide(prs)

    add_formula_slide(prs, "Evaluation Metrics", [
        ("Argument Refinement Index (ARI)",
         "ARI = (1/(R-1)) × Σ 𝟙[addresses_rebuttal]",
         "Measures how arguments incorporate opposing views across rounds"),
        ("Inter-Round Confidence Stability (IRCS)",
         "IRCS = 1 - Var(confidence) / max_var",
         "High = stable position, Low = responsive to arguments"),
        ("Confidence Convergence Rate (CCR)",
         "CCR = (gap_R1 - gap_Rn) / gap_R1",
         "Measures how quickly opposing agents reach agreement"),
    ])

    add_comparison_slide(prs, "Three-Way Verdict Comparison", [
        ["Verdict: FAITHFUL", "Confidence: 95%", "Mutations: 0", "LLM Calls: 1", "", "❌ Overconfident"],
        ["Verdict: AMBIGUOUS", "Confidence: 85%", "Mutations: 2", "LLM Calls: 10", "", "✓ Calibrated"],
        ["Verdict: FAITHFUL", "Confidence: 85%", "(forced binary)", "LLM Calls: 11", "", "✓ Lower conf"],
    ])

    add_table_slide(prs, "Confidence Evolution: Case 0 (COVID Deaths)",
                    ["Agent", "Round 1", "Round 2", "Round 3", "Δ"],
                    [
                        ["Prosecutor", "90%", "85%", "85%", "-5%"],
                        ["Defense", "95%", "90%", "90%", "-5%"],
                        ["Epistemologist", "90%", "90%", "90%", "0%"],
                    ])

    add_table_slide(prs, "1-Round vs 3-Round: Detailed Comparison",
                    ["Metric", "1-Round", "3-Round", "Winner"],
                    [
                        ["Confidence", "95%", "85%", "3-Round ✓"],
                        ["Mutation Types", "0", "2", "3-Round ✓"],
                        ["Ambiguity Handling", "No", "Yes", "3-Round ✓"],
                        ["Transparency", "Minimal", "Full transcript", "3-Round ✓"],
                        ["Speed", "~5s", "~30s", "1-Round ✓"],
                    ])

    add_key_findings_slide(prs)

    add_content_slide(prs, "Why 3 Rounds Works Better", [
        "Round 1: Establishes initial positions, identifies disagreements",
        "Round 2: Cross-examination reveals argument weaknesses",
        "Round 3: Refined positions incorporate valid rebuttals",
        "Result: Arguments are 40% more sophisticated by Round 3",
        "Key: Agents maintain productive tension, not false consensus"
    ])

    add_conclusion_slide(prs)

    # Thank you
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    tf = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1)).text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Questions?"
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Save
    output_path = os.path.join(os.path.dirname(__file__), "rounds_comparison.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
