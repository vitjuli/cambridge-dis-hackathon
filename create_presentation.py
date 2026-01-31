"""
Generate PowerPoint presentation for Multi-Agent vs Single-Agent analysis
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
import os

# Colors
DARK_BG = RGBColor(15, 23, 42)  # Slate 900
ACCENT = RGBColor(99, 102, 241)  # Indigo 500
SUCCESS = RGBColor(16, 185, 129)  # Emerald 500
WARNING = RGBColor(245, 158, 11)  # Amber 500
DANGER = RGBColor(239, 68, 68)  # Red 500
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(148, 163, 184)


def set_slide_background(slide, color):
    """Set solid background color for slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle):
    """Add title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Team badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(5.5), Inches(3), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(30, 27, 75)
    badge.line.color.rgb = ACCENT
    tf = badge.text_frame
    tf.paragraphs[0].text = "Kepler Team"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, title):
    """Add section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(30, 27, 75))

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.2), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    return slide


def add_content_slide(prs, title, content_items, has_metrics=False, metrics=None):
    """Add content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Accent line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(1.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # Content
    if has_metrics and metrics:
        # Add metrics boxes at top
        box_width = 2.8
        start_x = 0.5
        for i, (label, value, detail) in enumerate(metrics):
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(start_x + i * 3), Inches(1.5),
                                         Inches(box_width), Inches(1.3))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(30, 27, 75)
            box.line.color.rgb = RGBColor(55, 48, 107)

            tf = box.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(10)
            p.font.color.rgb = LIGHT_GRAY
            p.alignment = PP_ALIGN.CENTER

            p = tf.add_paragraph()
            p.text = value
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            p.alignment = PP_ALIGN.CENTER

            p = tf.add_paragraph()
            p.text = detail
            p.font.size = Pt(9)
            p.font.color.rgb = LIGHT_GRAY
            p.alignment = PP_ALIGN.CENTER

        content_start = 3.0
    else:
        content_start = 1.5

    # Bullet points
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(content_start), Inches(9), Inches(4))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)

    return slide


def add_comparison_slide(prs, title, single_data, multi_data):
    """Add comparison slide with two columns."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Single Agent Box
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(4.3), Inches(5))
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(30, 27, 75)
    box1.line.color.rgb = WARNING

    tf = box1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Single-Agent"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WARNING
    p.alignment = PP_ALIGN.CENTER

    for item in single_data:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_before = Pt(8)

    # Multi Agent Box
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.5), Inches(4.3), Inches(5))
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(30, 27, 75)
    box2.line.color.rgb = SUCCESS

    tf = box2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Multi-Agent Ensemble"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = SUCCESS
    p.alignment = PP_ALIGN.CENTER

    for item in multi_data:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_before = Pt(8)

    return slide


def add_table_slide(prs, title, headers, rows):
    """Add slide with table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Table
    cols = len(headers)
    table_rows = len(rows) + 1
    table = slide.shapes.add_table(table_rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5 * table_rows)).table

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(30, 27, 75)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

    return slide


def add_formula_slide(prs, title, formulas):
    """Add slide with formulas."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    y_pos = 1.5
    for name, formula, description in formulas:
        # Formula box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(9), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(30, 27, 75)
        box.line.color.rgb = RGBColor(55, 48, 107)

        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT

        p = tf.add_paragraph()
        p.text = formula
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.font.name = "Consolas"

        p = tf.add_paragraph()
        p.text = description
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_GRAY

        y_pos += 1.5

    return slide


def add_case_study_slide(prs, case_id, claim_short, single_verdict, multi_verdict, insight):
    """Add case study slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Case {case_id}: {claim_short}"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Verdicts comparison
    verdict_colors = {"FAITHFUL": SUCCESS, "MUTATED": DANGER, "AMBIGUOUS": WARNING}

    # Single verdict box
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(4.3), Inches(1.5))
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(30, 27, 75)
    box1.line.color.rgb = WARNING

    tf = box1.text_frame
    p = tf.paragraphs[0]
    p.text = "Single-Agent"
    p.font.size = Pt(12)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = single_verdict
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = verdict_colors.get(single_verdict, WHITE)
    p.alignment = PP_ALIGN.CENTER

    # Multi verdict box
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.5), Inches(4.3), Inches(1.5))
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(30, 27, 75)
    box2.line.color.rgb = SUCCESS

    tf = box2.text_frame
    p = tf.paragraphs[0]
    p.text = "Multi-Agent"
    p.font.size = Pt(12)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = multi_verdict
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = verdict_colors.get(multi_verdict, WHITE)
    p.alignment = PP_ALIGN.CENTER

    # Insight box
    insight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.3), Inches(9), Inches(3))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
    insight_box.line.color.rgb = ACCENT

    tf = insight_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key Insight"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    p = tf.add_paragraph()
    p.text = insight
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.space_before = Pt(8)

    return slide


def add_conclusion_slide(prs):
    """Add conclusion slide with key metrics."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(30, 27, 75))

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Conclusion"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Key metrics
    metrics = [
        ("50%", "ECE Improvement", "0.267 → 0.133"),
        ("2.67×", "Recall Improvement", "37.5% → 100%"),
        ("κ = 0.83", "Final Agreement", "Almost Perfect"),
    ]

    for i, (value, label, detail) in enumerate(metrics):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.8 + i * 3), Inches(2), Inches(2.6), Inches(2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(15, 23, 42)
        box.line.color.rgb = SUCCESS

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = SUCCESS
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = detail
        p.font.size = Pt(10)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

    # Key takeaway
    takeaway = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(9), Inches(1.8))
    takeaway.fill.solid()
    takeaway.fill.fore_color.rgb = RGBColor(15, 23, 42)
    takeaway.line.color.rgb = ACCENT

    tf = takeaway.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key Takeaway"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Multi-agent ensemble provides better calibrated confidence and transparent reasoning.\nFor fact-checking, knowing when to say \"uncertain\" is as valuable as being correct."
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide


def create_presentation():
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs,
                    "Multi-Agent Ensemble vs Single-Agent\nfor Claim Verification",
                    "A Statistical Analysis | Cambridge DIS Hackathon 2025")

    # Slide 2: Problem Statement
    add_content_slide(prs, "The Challenge", [
        "Fact verification requires distinguishing faithful representations from subtle mutations",
        "Single-agent systems often exhibit overconfidence and miss nuanced distortions",
        "Question: Can multi-agent debate improve verification quality?",
        "We compare: 1 LLM call (single) vs 13 LLM calls (4 agents × 4 rounds)"
    ])

    # Slide 3: Architecture
    add_content_slide(prs, "Multi-Agent Architecture", [
        "Prosecutor: Adversarial agent seeking evidence of mutation",
        "Defense: Charitable agent arguing for faithful interpretation",
        "Epistemologist: Meta-cognitive agent quantifying uncertainty",
        "Moderator: Synthesis agent aggregating final verdict",
        "4 debate rounds with confidence updates based on arguments"
    ])

    # Slide 4: Metrics
    add_section_slide(prs, "Evaluation Metrics")

    # Slide 5: Formulas
    add_formula_slide(prs, "Statistical Metrics", [
        ("Expected Calibration Error (ECE)",
         "ECE = Σ (|Bm|/n) × |acc(Bm) - conf(Bm)|",
         "Measures alignment between predicted confidence and actual accuracy"),
        ("Verdict Entropy",
         "H(V) = -Σ p(v) × log₂(p(v))",
         "Shannon entropy capturing uncertainty in system output"),
        ("Fleiss' Kappa",
         "κ = (P̄ - P̄e) / (1 - P̄e)",
         "Inter-rater agreement for multiple agents on categorical verdicts"),
    ])

    # Slide 6: Results Overview
    add_content_slide(prs, "Results Overview", [
        "ECE improved by 50%: 0.267 → 0.133 (better calibration)",
        "Verdict entropy: 0.918 → 1.585 bits (maximum = appropriate uncertainty)",
        "Mutation detection recall: 2.67× improvement (3 → 8 types detected)",
        "Inter-agent agreement: κ = 0.83 by final round (almost perfect)"
    ], has_metrics=True, metrics=[
        ("ECE", "0.133", "vs 0.267"),
        ("Entropy", "1.585", "bits (max)"),
        ("Recall", "2.67×", "improvement")
    ])

    # Slide 7: Comparison Table
    add_table_slide(prs, "Detailed Comparison",
                    ["Metric", "Single-Agent", "Multi-Agent", "Δ"],
                    [
                        ["Mean Confidence", "93.3%", "86.7%", "-6.6%"],
                        ["ECE", "0.267", "0.133", "-50%"],
                        ["Verdict Entropy", "0.918 bits", "1.585 bits", "+73%"],
                        ["Mutations Detected", "3", "8", "+167%"],
                        ["LLM Calls", "1", "12", "+11"],
                    ])

    # Slide 8: Calibration Analysis
    add_comparison_slide(prs, "Calibration Analysis",
                         [
                             "Confidence: 93.3%",
                             "Accuracy: 66.7%",
                             "ECE = |93.3 - 66.7| = 0.267",
                             "",
                             "❌ Overconfident",
                             "Case 0: 95% confident",
                             "but WRONG verdict"
                         ],
                         [
                             "Confidence: 86.7%",
                             "Accuracy: 66.7%",
                             "ECE = 0.133",
                             "",
                             "✓ Well-calibrated",
                             "Case 0: 80% confident",
                             "correctly marked AMBIGUOUS"
                         ])

    # Slide 9: Case Studies
    add_section_slide(prs, "Case Studies")

    # Slide 10: Case 0
    add_case_study_slide(prs, 0, "COVID Deaths - Numerical Boundary",
                         "FAITHFUL", "AMBIGUOUS",
                         "Single-agent missed the subtle shift from 'more than 14,500' (lower bound) "
                         "to 'less than 14,550' (upper bound). This framing change caps the death toll "
                         "rather than emphasizing the minimum. Multi-agent's Prosecutor caught this, "
                         "and after 4 rounds of debate, correctly concluded: genuinely ambiguous.")

    # Slide 11: Case 2
    add_case_study_slide(prs, 2, "Badmotorfinger - Sales Figures",
                         "MUTATED", "MUTATED",
                         "Both agreed on MUTATED, but multi-agent provided severity gradation: "
                         "'after March 1996' is LOW severity (acceptable paraphrase), while "
                         "'1.8 million sold' is HIGH severity (contradicts 1.5M in source). "
                         "Defense explicitly CONCEDED the numerical inflation point.")

    # Slide 12: Why Multi-Agent is Better
    add_content_slide(prs, "Why Multi-Agent Ensemble Wins", [
        "Better Calibrated: Lower confidence when uncertain (ECE -50%)",
        "Comprehensive Detection: 2.67× more mutation types found",
        "Transparent Reasoning: Full debate transcript for audit",
        "Handles Ambiguity: Uses 'uncertain' verdict appropriately",
        "Adversarial Testing: Agents attack each other's blind spots"
    ])

    # Slide 13: Agreement Evolution
    add_table_slide(prs, "Inter-Agent Agreement Evolution",
                    ["Round", "Mean κ", "Interpretation"],
                    [
                        ["1 (Initial)", "0.44", "Moderate agreement"],
                        ["2", "0.61", "Substantial agreement"],
                        ["3", "0.72", "Substantial agreement"],
                        ["4 (Final)", "0.83", "Almost perfect agreement"],
                    ])

    # Slide 14: Conclusion
    add_conclusion_slide(prs)

    # Slide 15: Thank You
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Questions?"
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Team info
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Kepler Team | github.com/Julia-Elisa/cambridge-dis-hackathon"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Save
    output_path = os.path.join(os.path.dirname(__file__), "kepler_presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
